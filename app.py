import sys
import sqlite3
import google.generativeai as genai
import speech_recognition as sr
import hashlib
import fitz  # PyMuPDF for PDF extraction
import pptx  # python-pptx for PPT extraction
from PIL import Image
import pytesseract
from PyQt5.QtWidgets import QApplication, QWidget, QVBoxLayout, QTextEdit, QPushButton, QLabel, QLineEdit, QMessageBox, QStackedWidget, QFileDialog
from PyQt5.QtGui import QPalette, QColor
from PyQt5.QtCore import Qt

# Set your Google Gemini AI API key
GEMINI_API_KEY = "your-api-key-here"

genai.configure(api_key=GEMINI_API_KEY)

# Hashing function for passwords
def hash_password(password):
    return hashlib.sha256(password.encode()).hexdigest()

# Database Setup
def init_db():
    conn = sqlite3.connect("notes.db")
    cursor = conn.cursor()
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT UNIQUE,
            password TEXT
        )
    """)
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS notes (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER,
            content TEXT,
            category TEXT,
            FOREIGN KEY(user_id) REFERENCES users(id)
        )
    """)
    conn.commit()
    conn.close()
init_db()

class LoginPage(QWidget):
    def __init__(self, stacked_widget):
        super().__init__()
        self.stacked_widget = stacked_widget
        self.setWindowTitle("Login")
        self.setGeometry(100, 100, 400, 200)
        self.setStyleSheet("background-color: #34495e; color: white;")
        layout = QVBoxLayout()

        self.label = QLabel("Login to your account")
        self.label.setAlignment(Qt.AlignCenter)
        self.label.setStyleSheet("font-size: 16px; font-weight: bold;")
        layout.addWidget(self.label)

        self.username_input = QLineEdit(self)
        self.username_input.setPlaceholderText("Username")
        layout.addWidget(self.username_input)

        self.password_input = QLineEdit(self)
        self.password_input.setPlaceholderText("Password")
        self.password_input.setEchoMode(QLineEdit.Password)
        layout.addWidget(self.password_input)

        self.login_button = QPushButton("Login", self)
        self.login_button.clicked.connect(self.authenticate_user)
        layout.addWidget(self.login_button)

        self.register_button = QPushButton("Register", self)
        self.register_button.clicked.connect(self.register_user)
        layout.addWidget(self.register_button)

        self.setLayout(layout)

    def authenticate_user(self):
        username = self.username_input.text().strip()
        password = hash_password(self.password_input.text().strip())
        
        conn = sqlite3.connect("notes.db")
        cursor = conn.cursor()
        cursor.execute("SELECT id FROM users WHERE username = ? AND password = ?", (username, password))
        user = cursor.fetchone()
        conn.close()

        if user:
            self.stacked_widget.setCurrentWidget(self.stacked_widget.widget(1))
        else:
            QMessageBox.warning(self, "Login Failed", "Invalid username or password.")

    def register_user(self):
        username = self.username_input.text().strip()
        password = hash_password(self.password_input.text().strip())
        
        if not username or not password:
            QMessageBox.warning(self, "Registration Failed", "Please enter a username and password.")
            return
        
        conn = sqlite3.connect("notes.db")
        cursor = conn.cursor()
        try:
            cursor.execute("INSERT INTO users (username, password) VALUES (?, ?)", (username, password))
            conn.commit()
            QMessageBox.information(self, "Success", "User registered successfully! You can now log in.")
        except sqlite3.IntegrityError:
            QMessageBox.warning(self, "Registration Failed", "Username already exists.")
        conn.close()

class AINoteSummarizer(QWidget):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("AI-Powered Voice Note Summarizer")
        self.setGeometry(100, 100, 600, 400)
        self.setStyleSheet("background-color: #2c3e50; color: white;")
        layout = QVBoxLayout()

        self.label = QLabel("Enter your note below or use voice input:")
        self.label.setAlignment(Qt.AlignCenter)
        self.label.setStyleSheet("font-size: 16px; font-weight: bold;")
        layout.addWidget(self.label)

        self.note_input = QTextEdit(self)
        self.note_input.setPlaceholderText("Write your note here or use voice input...")
        layout.addWidget(self.note_input)

        self.upload_button = QPushButton("Upload Document", self)
        self.upload_button.clicked.connect(self.upload_document)
        layout.addWidget(self.upload_button)

        self.summarize_button = QPushButton("Summarize", self)
        self.summarize_button.clicked.connect(self.summarize_note)
        layout.addWidget(self.summarize_button)

        self.export_button = QPushButton("Export as PDF", self)
        self.export_button.clicked.connect(self.export_to_pdf)
        layout.addWidget(self.export_button)

        self.summary_output = QTextEdit(self)
        self.summary_output.setReadOnly(True)
        layout.addWidget(self.summary_output)

        self.setLayout(layout)

    def upload_document(self):
        file_path, _ = QFileDialog.getOpenFileName(self, "Open File", "", "Documents (*.pdf *.pptx *.png *.jpg)")
        if file_path:
            extracted_text = self.extract_text_from_file(file_path)
            self.note_input.setPlainText(extracted_text)
    
    def extract_text_from_file(self, file_path):
        if file_path.endswith(".pdf"):
            doc = fitz.open(file_path)
            text = "\n".join([page.get_text("text") for page in doc])
            return text
        elif file_path.endswith(".pptx"):
            prs = pptx.Presentation(file_path)
            text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text')])
            return text
        elif file_path.endswith(".png") or file_path.endswith(".jpg"):
            img = Image.open(file_path)
            return pytesseract.image_to_string(img)
        return "Unsupported file type."

    def summarize_note(self):
        note_text = self.note_input.toPlainText().strip()
        
        if not note_text:
            self.summary_output.setPlainText("Please enter a note to summarize.")
            return

        try:
            print("Generating summary...")  # Debugging Step 1

            model = genai.GenerativeModel("gemini-pro")
            response = model.generate_content(note_text)

            print("API Response:", response)  # Debugging Step 2

            if hasattr(response, "text"):
                summary = response.text.strip()
                print("Generated Summary:", summary)  # Debugging Step 3
            else:
                summary = "Error: No summary generated."

        except Exception as e:
            summary = f"Error: {str(e)}"
            print("Error Occurred:", summary)  # Debugging Step 4

        self.summary_output.setPlainText(summary)

    def export_to_pdf(self):
        file_path, _ = QFileDialog.getSaveFileName(self, "Save as PDF", "", "PDF Files (*.pdf)")
        if file_path:
            doc = fitz.open()
            page = doc.new_page()
            page.insert_text((50, 100), self.summary_output.toPlainText())
            doc.save(file_path)
            doc.close()
            QMessageBox.information(self, "Success", "Summary exported as PDF successfully!")

if __name__ == "__main__":
    app = QApplication(sys.argv)
    stacked_widget = QStackedWidget()
    login_page = LoginPage(stacked_widget)
    summarizer_page = AINoteSummarizer()
    stacked_widget.addWidget(login_page)
    stacked_widget.addWidget(summarizer_page)
    stacked_widget.setCurrentWidget(login_page)
    stacked_widget.show()
    sys.exit(app.exec_())
